"""Multi-format document text extraction and chunking for digital twins."""
from __future__ import annotations

import csv as _csv
import hashlib
import logging
import html
import json
import mimetypes
import os
import re
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

try:
    import tomllib as _tomllib  # type: ignore
except Exception:
    _tomllib = None

LOGGER = logging.getLogger(__name__)

TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".r", ".sh", ".json", ".yaml", ".yml", ".sql", ".csv", ".tsv",
    ".html", ".htm", ".xml", ".svg", ".drawio", ".ipynb", ".rtf",
    ".toml", ".ini", ".cfg", ".log", ".jsonl", ".ts", ".tsx", ".jsx",
    ".java", ".c", ".cpp", ".h", ".hpp", ".go", ".rs", ".rb", ".pl",
    ".scala", ".jl", ".m", ".css", ".scss", ".less", ".js", ".dart", ".vue", ".svelte",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg", ".tif", ".tiff", ".bmp"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".doc", ".dotx", ".odt", ".rtf"}
PRESENTATION_EXTENSIONS = {".pptx", ".ppt", ".odp"}
DATA_EXTENSIONS = {".xlsx", ".xls", ".csv", ".tsv", ".ods", ".parquet", ".feather", ".h5", ".hdf5", ".rds", ".rdata"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".avi", ".webm"}
SCANNABLE_EXTENSIONS = (
    TEXT_EXTENSIONS | IMAGE_EXTENSIONS | DOCUMENT_EXTENSIONS
    | PRESENTATION_EXTENSIONS | DATA_EXTENSIONS | VIDEO_EXTENSIONS | {".drawio"}
)
EXTRACTABLE_TEXT_EXTENSIONS = (
    TEXT_EXTENSIONS | DOCUMENT_EXTENSIONS | PRESENTATION_EXTENSIONS | DATA_EXTENSIONS
)
BINARY_METADATA_ONLY_EXTENSIONS = IMAGE_EXTENSIONS | VIDEO_EXTENSIONS | {".parquet", ".feather", ".h5", ".hdf5", ".rds", ".rdata"}
SKIP_PARTS = {".git", "node_modules", ".venv", "__pycache__", "media", ".DS_Store", ".dart_tool", "build", "dist", ".next", ".pytest_cache", ".mypy_cache", ".ruff_cache", "coverage"}

LIFECYCLE_RULES: list[tuple[str, str, str]] = [
    ("management", "Management & Planning", r"management|planning"),
    ("methods", "Methods & Experiments", r"method|experiment|protocol|wet.?lab|cycif|geomx|xenium"),
    ("data_figures", "Data & Figures", r"data.*figure|data & figure|data and figure|\bfigures?\b|\bdata/|\bdata\\"),
    ("meetings", "Meetings & Updates", r"meeting|update"),
    ("writing", "Writing & Dissemination", r"writing|dissemination|manuscript|abstract|poster|final_manuscript"),
    ("presentations", "Presentations & Slides", r"presentation|slide|poster"),
    ("archive", "Archive", r"archive"),
    ("root", "Project Root", r"^$"),
]
FIGURE_PATH_HINTS = re.compile(
    r"figure|fig[\d_\-.]|snapshot|screenshot|plot|spatial|tsi|rcn|thesis|finished.?fig|images and fig|supplementary.?fig",
    re.I,
)
FIGURE_NAME_HINTS = re.compile(r"^fig[\d_\-.]|^s\d+_p|\.png$|\.svg$|\.jpeg$|\.jpg$", re.I)
DOI_PATTERN = re.compile(r"10\.\d{4,9}/[^\s\])>]+", re.I)
TODO_ITEM = re.compile(r"(?:^|\n)\s*(?:[-*]|\d+\.)\s*(?:TODO|To do|Action|Next step)[:\s]+(.+)", re.I)


def _clean(text: str) -> str:
    text = re.sub(r"\[\[([^\]]+)\]\]\([^)]+\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"\{[^}]+\}", "", text)
    text = re.sub(r"\*{1,2}([^*]+)\*{1,2}", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r">\s*", "", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _classify_lifecycle_section(rel_path: str) -> str:
    parts = Path(rel_path).parts
    probe = " ".join(parts[:3]).lower() if parts else ""
    if len(parts) <= 1 and (not parts or not parts[0].startswith(("1", "2", "3", "4", "5", "6"))):
        return "root"
    for sec_id, _label, pattern in LIFECYCLE_RULES:
        if sec_id == "root":
            continue
        if re.search(pattern, probe, re.I):
            return sec_id
    return "root"


def _classify_asset_type(ext: str, rel_path: str, name: str) -> str:
    p = rel_path.lower()
    n = name.lower()
    if ext in IMAGE_EXTENSIONS:
        if FIGURE_PATH_HINTS.search(p) or FIGURE_NAME_HINTS.search(n) or "figure" in p:
            return "figure"
        return "image"
    if ext in DOCUMENT_EXTENSIONS:
        if "fig" in n and ext == ".pdf":
            return "figure"
        return "document"
    if ext in PRESENTATION_EXTENSIONS:
        return "presentation"
    if ext in DATA_EXTENSIONS:
        return "data"
    if ext in TEXT_EXTENSIONS:
        return "text"
    if ext in VIDEO_EXTENSIONS:
        return "video"
    return "other"


def _figure_priority(rel_path: str, name: str) -> int:
    p, n = rel_path.lower(), name.lower()
    score = 0
    if re.match(r"^figures?/", p) or p.startswith("figures/"):
        score += 12
    if "final_manuscript" in p:
        score += 10
    if re.search(r"fig[\d_\-.]", n):
        score += 9
    if FIGURE_PATH_HINTS.search(p):
        score += 5
    if "screenshot" in p or "snapshot" in p or "spatial_count" in p:
        score -= 4
    if "archive" in p:
        score += 2
    depth = len(Path(rel_path).parts)
    if depth > 7:
        score -= min(depth - 7, 8)
    if n.endswith(".svg"):
        score += 1
    return score


DEFAULT_MAX_TEXT_CHARS_PER_FILE = int(os.environ.get("DIGITAL_TWIN_MAX_TEXT_CHARS_PER_FILE", "250000"))
DEFAULT_CHUNK_CHARS = int(os.environ.get("DIGITAL_TWIN_CHUNK_CHARS", "1800"))
DEFAULT_CHUNK_OVERLAP = int(os.environ.get("DIGITAL_TWIN_CHUNK_OVERLAP", "250"))
DEFAULT_MAX_PDF_PAGES = int(os.environ.get("DIGITAL_TWIN_MAX_PDF_PAGES", "120"))
DEFAULT_MAX_PPTX_SLIDES = int(os.environ.get("DIGITAL_TWIN_MAX_PPTX_SLIDES", "180"))
DEFAULT_MAX_XLSX_SHEETS = int(os.environ.get("DIGITAL_TWIN_MAX_XLSX_SHEETS", "20"))
DEFAULT_MAX_XLSX_ROWS = int(os.environ.get("DIGITAL_TWIN_MAX_XLSX_ROWS", "200"))
DEFAULT_MAX_CSV_ROWS = int(os.environ.get("DIGITAL_TWIN_MAX_CSV_ROWS", "300"))
DEFAULT_MAX_DOCS_IN_JSON = int(os.environ.get("DIGITAL_TWIN_MAX_DOCS_IN_JSON", "800"))
DEFAULT_MAX_CHUNKS_IN_JSON = int(os.environ.get("DIGITAL_TWIN_MAX_CHUNKS_IN_JSON", "2500"))
DEFAULT_SUBPROCESS_TIMEOUT = int(os.environ.get("DIGITAL_TWIN_EXTRACT_TIMEOUT", "40"))
DEFAULT_MAX_FILE_BYTES = int(os.environ.get("DIGITAL_TWIN_MAX_FILE_BYTES", str(512 * 1024 * 1024)))

DOCUMENT_KIND_BY_EXTENSION = {
    ".pdf": "pdf",
    ".docx": "word_document",
    ".doc": "legacy_word_document",
    ".dotx": "word_template",
    ".odt": "open_document_text",
    ".rtf": "rich_text",
    ".pptx": "presentation",
    ".ppt": "legacy_presentation",
    ".odp": "open_document_presentation",
    ".xlsx": "spreadsheet",
    ".xls": "legacy_spreadsheet",
    ".ods": "open_document_spreadsheet",
    ".csv": "table_csv",
    ".tsv": "table_tsv",
    ".ipynb": "notebook",
    ".drawio": "diagram",
    ".svg": "vector_image",
    ".html": "html",
    ".htm": "html",
    ".xml": "xml",
    ".json": "json",
    ".jsonl": "jsonl",
    ".yaml": "yaml",
    ".yml": "yaml",
    ".toml": "toml",
    ".md": "markdown",
    ".txt": "plain_text",
    ".py": "python_script",
    ".r": "r_script",
    ".sh": "shell_script",
    ".sql": "sql_script",
}


@dataclass
class ExtractionResult:
    path: str
    name: str
    extension: str
    document_kind: str
    mime_type: str | None
    size_bytes: int
    modified_at: str | None
    sha256: str | None = None
    status: str = "not_attempted"  # extracted | metadata_only | skipped | failed | empty
    extractor: str = "none"
    text: str = ""
    excerpt: str = ""
    title: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)
    char_count: int = 0
    word_count: int = 0
    chunks: list[dict[str, Any]] = field(default_factory=list)

    def as_json(self, include_text: bool = False, include_chunks: bool = True) -> dict[str, Any]:
        payload = {
            "path": self.path,
            "name": self.name,
            "extension": self.extension,
            "document_kind": self.document_kind,
            "mime_type": self.mime_type,
            "size_bytes": self.size_bytes,
            "modified_at": self.modified_at,
            "sha256": self.sha256,
            "status": self.status,
            "extractor": self.extractor,
            "title": self.title,
            "excerpt": self.excerpt,
            "char_count": self.char_count,
            "word_count": self.word_count,
            "metadata": self.metadata,
            "warnings": self.warnings,
            "errors": self.errors,
        }
        if include_text:
            payload["text"] = self.text
        if include_chunks:
            payload["chunks"] = self.chunks
        return payload


def _utc_iso_from_timestamp(ts: float | int | None) -> str | None:
    if ts is None:
        return None
    try:
        return datetime.fromtimestamp(float(ts), timezone.utc).isoformat()
    except Exception:
        return None


def _safe_stat(path: Path) -> dict[str, Any]:
    try:
        st = path.stat()
        return {
            "size_bytes": int(st.st_size),
            "modified_at": _utc_iso_from_timestamp(st.st_mtime),
        }
    except OSError:
        return {"size_bytes": 0, "modified_at": None}


def _sha256_file(path: Path, limit_bytes: int | None = None) -> str | None:
    try:
        h = hashlib.sha256()
        read = 0
        with path.open("rb") as fh:
            for block in iter(lambda: fh.read(1024 * 1024), b""):
                if limit_bytes is not None and read + len(block) > limit_bytes:
                    block = block[: max(0, limit_bytes - read)]
                if not block:
                    break
                h.update(block)
                read += len(block)
                if limit_bytes is not None and read >= limit_bytes:
                    break
        return h.hexdigest()
    except OSError:
        return None


def _limit_text(text: str, max_chars: int = DEFAULT_MAX_TEXT_CHARS_PER_FILE) -> tuple[str, bool]:
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars], True


def _normalize_extracted_text(text: str) -> str:
    if not text:
        return ""
    text = html.unescape(text)
    text = text.replace("\x00", " ")
    text = re.sub(r"[\u200b\u200c\u200d\ufeff]", "", text)
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    return text.strip()


def _clean_for_excerpt(text: str, max_chars: int = 800) -> str:
    """Prefer the first substantive paragraph, not page/slide/XML scaffolding."""
    lines: list[str] = []
    for raw in (text or "").splitlines()[:120]:
        line = _clean(raw)
        if len(line) < 24 or _is_junk_title_line(line):
            continue
        lines.append(line)
        if sum(len(x) + 1 for x in lines) >= max_chars * 2:
            break
    if lines:
        return " ".join(lines)[:max_chars]
    return _clean(text[: max_chars * 2])[:max_chars]


def _humanize_filename(name: str) -> str:
    stem = Path(name).stem
    stem = stem.replace("_", " ").replace("-", " ")
    stem = re.sub(r"\s+", " ", stem).strip()
    return stem[:160] if stem else "Document"


_JUNK_TITLE_RE = re.compile(
    r"^(?:#+\s*)?(?:page\s*\d+(?:\s+of\s+\d+)?|slide\s*\d+|ppt/slides/|word/document|\d{1,2}\.\d{1,2}\.\d{4})$",
    re.I,
)


def _is_junk_title_line(line: str) -> bool:
    s = (line or "").strip()
    if not s or len(s) < 8:
        return True
    if _JUNK_TITLE_RE.match(s):
        return True
    if s.startswith("### ") and ("/" in s or ".xml" in s.lower()):
        return True
    if s.startswith("## Page ") or s.startswith("## Slide "):
        return True
    return False


def document_display_title(doc: dict[str, Any]) -> str:
    """Human-readable label for UI lists (filename-first, not raw extractor headings)."""
    path = (doc.get("path") or doc.get("relative_path") or "").replace("\\", "/")
    name = doc.get("name") or (Path(path).name if path else "")
    human = _humanize_filename(name) if name else "Document"
    raw = _clean((doc.get("title") or "").strip())
    if raw and not _is_junk_title_line(raw) and raw.lower() != human.lower():
        if len(raw) <= 120:
            return raw
    return human


def document_display_excerpt(doc: dict[str, Any], max_chars: int = 500) -> str | None:
    excerpt = (doc.get("excerpt") or "").strip()
    if excerpt:
        cleaned = _clean_for_excerpt(excerpt, max_chars=max_chars)
        if cleaned and not _is_junk_title_line(cleaned):
            return cleaned
    return None


def lab_database_asset_url(relative_root: str, relative_path: str) -> str:
    from urllib.parse import quote

    root = (relative_root or "").replace("\\", "/").strip("/")
    file_part = (relative_path or "").replace("\\", "/").lstrip("/")
    combined = "/".join(p for p in (root, file_part) if p)
    return "/database-static/" + "/".join(quote(seg, safe="") for seg in combined.split("/"))


def _count_words(text: str) -> int:
    return len(re.findall(r"[A-Za-z0-9\u00C0-\uFFFF]+", text or ""))


def _chunk_text(text: str, source_path: str, *, chunk_chars: int = DEFAULT_CHUNK_CHARS,
                overlap: int = DEFAULT_CHUNK_OVERLAP) -> list[dict[str, Any]]:
    text = _normalize_extracted_text(text)
    if not text:
        return []
    try:
        from app_skeleton.api.chunking import chunk_text as facade_chunk_text
        from app_skeleton.api.platform_flags import canonical_chunk_pipeline_enabled

        chunks = facade_chunk_text(text, section_path=source_path)
        if chunks:
            return chunks
        if canonical_chunk_pipeline_enabled():
            return []
    except Exception:
        from app_skeleton.api.platform_flags import canonical_chunk_pipeline_enabled

        if canonical_chunk_pipeline_enabled():
            LOGGER.warning("canonical chunk pipeline failed for %s", source_path, exc_info=True)
            return []
        LOGGER.debug("chunking facade unavailable; using legacy _chunk_text", exc_info=True)
    chunk_chars = max(500, int(chunk_chars))
    overlap = min(max(0, int(overlap)), chunk_chars // 2)
    chunks: list[dict[str, Any]] = []
    start = 0
    idx = 0
    while start < len(text):
        end = min(len(text), start + chunk_chars)
        if end < len(text):
            # Prefer ending at paragraph/sentence/line boundary near the end.
            window = text[start:end]
            candidates = [window.rfind("\n\n"), window.rfind(". "), window.rfind("\n"), window.rfind("; ")]
            boundary = max(candidates)
            if boundary > int(chunk_chars * 0.55):
                end = start + boundary + 1
        chunk = text[start:end].strip()
        if chunk:
            chunks.append({
                "chunk_id": f"{source_path}::chunk_{idx:04d}",
                "source_file": source_path,
                "chunk_index": idx,
                "start_char": start,
                "end_char": end,
                "char_count": len(chunk),
                "word_count": _count_words(chunk),
                "text": chunk,
            })
            idx += 1
        if end >= len(text):
            break
        start = max(end - overlap, start + 1)
    return chunks


def _command_exists(name: str) -> bool:
    return shutil.which(name) is not None


def _run_command(args: list[str], *, timeout: int = DEFAULT_SUBPROCESS_TIMEOUT) -> tuple[int, str, str]:
    """Run an extractor command with bounded runtime and sanitized environment."""
    if not args:
        return 1, "", "empty command"
    try:
        env = {k: v for k, v in os.environ.items() if k not in {"PYTHONPATH"}}
        env.setdefault("LC_ALL", "C.UTF-8")
        proc = subprocess.run(
            args,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            timeout=max(1, int(timeout)),
            check=False,
            env=env,
        )
        return proc.returncode, proc.stdout or "", proc.stderr or ""
    except subprocess.TimeoutExpired as e:
        return 124, e.stdout or "", f"timeout after {timeout}s"
    except Exception as e:  # pragma: no cover
        LOGGER.debug("Extractor command failed: %s", e)
        return 1, "", str(e)


def _read_text_file(path: Path, max_chars: int = DEFAULT_MAX_TEXT_CHARS_PER_FILE) -> tuple[str, str, list[str]]:
    warnings: list[str] = []
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            text = path.read_text(encoding=enc, errors="replace")
            limited, truncated = _limit_text(text, max_chars)
            if truncated:
                warnings.append(f"truncated to {max_chars} characters")
            return _normalize_extracted_text(limited), f"text:{enc}", warnings
        except UnicodeDecodeError:
            continue
        except OSError as e:
            return "", "text", [str(e)]
    return "", "text", ["could not decode text file"]


def _strip_html_text(text: str) -> str:
    try:
        from bs4 import BeautifulSoup  # type: ignore
        soup = BeautifulSoup(text, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.extract()
        return soup.get_text("\n")
    except Exception:
        text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", text)
        text = re.sub(r"(?s)<[^>]+>", " ", text)
        return text


def _extract_json_text(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
        if path.suffix.lower() == ".jsonl":
            lines = []
            for i, line in enumerate(raw.splitlines()[:DEFAULT_MAX_CSV_ROWS]):
                line = line.strip()
                if not line:
                    continue
                try:
                    obj = json.loads(line)
                    lines.append(json.dumps(obj, ensure_ascii=False, indent=2))
                except Exception:
                    lines.append(line)
            return "\n\n".join(lines), "jsonl", {"line_limit": DEFAULT_MAX_CSV_ROWS}, warnings
        obj = json.loads(raw)
        return json.dumps(obj, ensure_ascii=False, indent=2), "json", {"top_level_type": type(obj).__name__}, warnings
    except Exception as e:
        text, extractor, ws = _read_text_file(path)
        return text, extractor, {}, warnings + ws + [f"json parse failed: {e}"]


def _extract_ipynb_text(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    try:
        nb = json.loads(path.read_text(encoding="utf-8", errors="replace"))
        lines: list[str] = []
        code_cells = 0
        markdown_cells = 0
        for idx, cell in enumerate(nb.get("cells", [])):
            cell_type = cell.get("cell_type", "cell")
            source = cell.get("source", "")
            if isinstance(source, list):
                source = "".join(source)
            if not source:
                continue
            if cell_type == "code":
                code_cells += 1
            if cell_type == "markdown":
                markdown_cells += 1
            lines.append(f"## Cell {idx + 1} ({cell_type})\n{source}")
        metadata = {
            "cell_count": len(nb.get("cells", [])),
            "code_cells": code_cells,
            "markdown_cells": markdown_cells,
            "kernel": (nb.get("metadata", {}).get("kernelspec") or {}).get("display_name"),
        }
        return "\n\n".join(lines), "ipynb-json", metadata, warnings
    except Exception as e:
        return "", "ipynb-json", {}, [str(e)]


def _xml_text_from_zip(path: Path, member_patterns: tuple[str, ...]) -> tuple[str, list[str], dict[str, Any]]:
    warnings: list[str] = []
    metadata: dict[str, Any] = {"zip_members_read": 0}
    texts: list[str] = []
    try:
        with zipfile.ZipFile(path) as zf:
            names = zf.namelist()
            selected = []
            for name in names:
                lower = name.lower()
                if any(re.search(pattern, lower) for pattern in member_patterns):
                    selected.append(name)
            for name in sorted(selected):
                try:
                    raw = zf.read(name)
                    metadata["zip_members_read"] += 1
                    root = ET.fromstring(raw)
                    chunks = [node.text for node in root.iter() if node.text and node.text.strip()]
                    if chunks:
                        texts.append(f"\n### {name}\n" + "\n".join(chunks))
                except Exception as e:
                    warnings.append(f"could not parse {name}: {e}")
        return "\n".join(texts), warnings, metadata
    except Exception as e:
        return "", [str(e)], metadata


def _extract_docx_text(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    metadata: dict[str, Any] = {}
    try:
        from docx import Document  # type: ignore
        doc = Document(str(path))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text and para.text.strip():
                parts.append(para.text)
        table_count = 0
        for table in doc.tables:
            table_count += 1
            for row in table.rows:
                vals = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if vals:
                    parts.append(" | ".join(vals))
        metadata["paragraph_count"] = len(doc.paragraphs)
        metadata["table_count"] = table_count
        try:
            cp = doc.core_properties
            metadata.update({
                "author": cp.author,
                "title": cp.title,
                "subject": cp.subject,
                "created": cp.created.isoformat() if cp.created else None,
                "modified": cp.modified.isoformat() if cp.modified else None,
            })
        except Exception:
            pass
        return "\n".join(parts), "python-docx", metadata, warnings
    except Exception as e:
        warnings.append(f"python-docx unavailable/failed: {e}")
        text, ws, md = _xml_text_from_zip(path, (r"word/document\.xml$", r"word/header\d*\.xml$", r"word/footer\d*\.xml$", r"word/footnotes\.xml$", r"word/endnotes\.xml$"))
        metadata.update(md)
        return text, "docx-zip-xml", metadata, warnings + ws


def _extract_odt_text(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    text, warnings, md = _xml_text_from_zip(path, (r"content\.xml$", r"meta\.xml$"))
    return text, "odf-zip-xml", md, warnings


def _extract_pptx_slide_xml(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    """Zip/XML fallback with slide numbers (no raw member paths in text)."""
    warnings: list[str] = []
    metadata: dict[str, Any] = {"slide_count": 0}
    parts: list[str] = []
    try:
        with zipfile.ZipFile(path) as zf:
            slide_names = sorted(
                [n for n in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n, re.I)],
                key=lambda n: int(re.search(r"slide(\d+)", n, re.I).group(1)),
            )
            metadata["slide_count"] = len(slide_names)
            for name in slide_names[:DEFAULT_MAX_PPTX_SLIDES]:
                slide_no = int(re.search(r"slide(\d+)", name, re.I).group(1))
                try:
                    root = ET.fromstring(zf.read(name))
                    chunks = [node.text.strip() for node in root.iter() if node.text and node.text.strip()]
                    if chunks:
                        parts.append(f"## Slide {slide_no}\n" + "\n".join(chunks))
                except Exception as e:
                    warnings.append(f"could not parse {name}: {e}")
            if len(slide_names) > DEFAULT_MAX_PPTX_SLIDES:
                warnings.append(f"read first {DEFAULT_MAX_PPTX_SLIDES} slides of {len(slide_names)}")
        return "\n\n".join(parts), "pptx-slide-xml", metadata, warnings
    except Exception as e:
        return "", "pptx-slide-xml", metadata, [str(e)]


def _extract_pptx_text(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    metadata: dict[str, Any] = {}
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(path))
        parts: list[str] = []
        slide_count = len(prs.slides)
        metadata["slide_count"] = slide_count
        max_slides = min(slide_count, DEFAULT_MAX_PPTX_SLIDES)
        # Index access — slicing prs.slides can break on some ONCOSYS decks (rId list bug).
        for idx in range(max_slides):
            slide_no = idx + 1
            try:
                slide = prs.slides[idx]
            except Exception as e:
                warnings.append(f"slide {slide_no} unavailable: {e}")
                continue
            slide_lines: list[str] = []
            try:
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        slide_lines.append(text.strip())
                    if getattr(shape, "has_table", False):
                        try:
                            for row in shape.table.rows:
                                vals = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                                if vals:
                                    slide_lines.append(" | ".join(vals))
                        except Exception:
                            pass
            except Exception as e:
                warnings.append(f"slide {slide_no} shape read failed: {e}")
                continue
            try:
                notes = slide.notes_slide.notes_text_frame.text
                if notes and notes.strip():
                    slide_lines.append("Speaker notes: " + notes.strip())
            except Exception:
                pass
            if slide_lines:
                parts.append(f"## Slide {slide_no}\n" + "\n".join(slide_lines))
        if slide_count > DEFAULT_MAX_PPTX_SLIDES:
            warnings.append(f"read first {DEFAULT_MAX_PPTX_SLIDES} slides of {slide_count}")
        if parts:
            return "\n\n".join(parts), "python-pptx", metadata, warnings
        warnings.append("python-pptx returned no slide text; using zip/XML fallback")
    except Exception as e:
        warnings.append(f"python-pptx unavailable/failed: {e}")
    text, extractor, md, ws = _extract_pptx_slide_xml(path)
    metadata.update(md)
    return text, extractor, metadata, warnings + ws


def _extract_pdf_text(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    metadata: dict[str, Any] = {}
    # Best lightweight option first.
    for module_name in ("pypdf", "PyPDF2"):
        try:
            mod = __import__(module_name)
            PdfReader = getattr(mod, "PdfReader")
            reader = PdfReader(str(path))
            pages = reader.pages
            metadata["page_count"] = len(pages)
            try:
                doc_md = getattr(reader, "metadata", None) or {}
                for k, v in dict(doc_md).items():
                    metadata[str(k).strip("/")] = str(v) if v is not None else None
            except Exception:
                pass
            text_parts = []
            max_pages = min(len(pages), DEFAULT_MAX_PDF_PAGES)
            for i in range(max_pages):
                try:
                    page_text = pages[i].extract_text() or ""
                    if page_text.strip():
                        text_parts.append(f"## Page {i + 1}\n{page_text}")
                except Exception as e:
                    warnings.append(f"page {i + 1} extraction failed: {e}")
            if len(pages) > DEFAULT_MAX_PDF_PAGES:
                warnings.append(f"read first {DEFAULT_MAX_PDF_PAGES} pages of {len(pages)}")
            return "\n\n".join(text_parts), module_name, metadata, warnings
        except Exception as e:
            warnings.append(f"{module_name} unavailable/failed: {e}")
    # CLI fallback: pdftotext is common on Linux/macOS.
    if _command_exists("pdftotext"):
        code, out, err = _run_command(["pdftotext", "-layout", "-nopgbrk", str(path), "-"])
        if code == 0 and out.strip():
            return out, "pdftotext", metadata, warnings
        warnings.append(f"pdftotext failed: {err.strip()[:300]}")
    return "", "pdf", metadata, warnings


def _extract_xlsx_text(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    metadata: dict[str, Any] = {}
    try:
        from openpyxl import load_workbook  # type: ignore
        wb = load_workbook(filename=str(path), read_only=True, data_only=True)
        sheet_names = wb.sheetnames
        metadata["sheet_names"] = sheet_names
        metadata["sheet_count"] = len(sheet_names)
        parts: list[str] = []
        for sheet_name in sheet_names[:DEFAULT_MAX_XLSX_SHEETS]:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            row_count = 0
            for row in ws.iter_rows(max_row=DEFAULT_MAX_XLSX_ROWS, values_only=True):
                vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if vals:
                    parts.append("\t".join(vals))
                row_count += 1
            if ws.max_row and ws.max_row > DEFAULT_MAX_XLSX_ROWS:
                warnings.append(f"sheet {sheet_name}: read first {DEFAULT_MAX_XLSX_ROWS} rows of {ws.max_row}")
        if len(sheet_names) > DEFAULT_MAX_XLSX_SHEETS:
            warnings.append(f"read first {DEFAULT_MAX_XLSX_SHEETS} sheets of {len(sheet_names)}")
        try:
            wb.close()
        except Exception:
            pass
        return "\n".join(parts), "openpyxl", metadata, warnings
    except Exception as e:
        warnings.append(f"openpyxl unavailable/failed: {e}")
        text, ws, md = _xml_text_from_zip(path, (r"xl/sharedStrings\.xml$", r"xl/worksheets/sheet\d+\.xml$"))
        metadata.update(md)
        return text, "xlsx-zip-xml", metadata, warnings + ws


def _extract_csv_text(path: Path, delimiter: str | None = None) -> tuple[str, str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    metadata: dict[str, Any] = {}
    try:
        sample = path.read_text(encoding="utf-8-sig", errors="replace")[:8192]
        if delimiter is None:
            try:
                dialect = _csv.Sniffer().sniff(sample)
                delimiter = dialect.delimiter
            except Exception:
                delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
        rows: list[list[str]] = []
        with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as fh:
            reader = _csv.reader(fh, delimiter=delimiter)
            for i, row in enumerate(reader):
                if i >= DEFAULT_MAX_CSV_ROWS:
                    warnings.append(f"read first {DEFAULT_MAX_CSV_ROWS} rows")
                    break
                rows.append([str(v).strip() for v in row])
        metadata["rows_read"] = len(rows)
        metadata["delimiter"] = delimiter
        metadata["column_count_guess"] = max((len(r) for r in rows), default=0)
        text_lines = ["\t".join([v for v in row if v]) for row in rows]
        return "\n".join(text_lines), "csv", metadata, warnings
    except Exception as e:
        text, extractor, ws = _read_text_file(path)
        return text, extractor, metadata, warnings + ws + [str(e)]


def _extract_xmlish_text(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
        # DrawIO files often contain diagrams as XML with encoded labels.
        if path.suffix.lower() in {".html", ".htm"}:
            return _strip_html_text(raw), "html", {}, warnings
        try:
            root = ET.fromstring(raw)
            texts = []
            for node in root.iter():
                if node.text and node.text.strip():
                    texts.append(node.text.strip())
                for attr in ("label", "value", "title", "name", "content"):
                    val = node.attrib.get(attr)
                    if val and val.strip():
                        texts.append(val.strip())
            return "\n".join(texts) or raw, "xml", {"root_tag": root.tag}, warnings
        except Exception:
            return re.sub(r"<[^>]+>", " ", raw), "xml-regex", {}, warnings
    except Exception as e:
        return "", "xml", {}, [str(e)]


def _extract_rtf_text(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    if _command_exists("textutil"):
        code, out, err = _run_command(["textutil", "-convert", "txt", "-stdout", str(path)])
        if code == 0 and out.strip():
            return out, "textutil", {}, warnings
        warnings.append(f"textutil failed: {err.strip()[:300]}")
    text, extractor, ws = _read_text_file(path)
    # Very rough RTF fallback.
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
    text = re.sub(r"[{}]", " ", text)
    return text, extractor + ":rtf-strip", {}, warnings + ws


def _extract_legacy_office_via_cli(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    ext = path.suffix.lower()
    if _command_exists("textutil") and ext in {".doc", ".rtf"}:
        code, out, err = _run_command(["textutil", "-convert", "txt", "-stdout", str(path)])
        if code == 0 and out.strip():
            return out, "textutil", {}, warnings
        warnings.append(f"textutil failed: {err.strip()[:300]}")
    if ext == ".doc" and _command_exists("antiword"):
        code, out, err = _run_command(["antiword", str(path)])
        if code == 0 and out.strip():
            return out, "antiword", {}, warnings
        warnings.append(f"antiword failed: {err.strip()[:300]}")
    if ext == ".doc" and _command_exists("catdoc"):
        code, out, err = _run_command(["catdoc", str(path)])
        if code == 0 and out.strip():
            return out, "catdoc", {}, warnings
        warnings.append(f"catdoc failed: {err.strip()[:300]}")
    if ext in {".xls", ".ods"} and _command_exists("in2csv"):
        code, out, err = _run_command(["in2csv", str(path)])
        if code == 0 and out.strip():
            return out, "in2csv", {}, warnings
        warnings.append(f"in2csv failed: {err.strip()[:300]}")
    if _command_exists("libreoffice") or _command_exists("soffice"):
        exe = shutil.which("libreoffice") or shutil.which("soffice")
        with tempfile.TemporaryDirectory(prefix="dt_extract_") as tmp:
            # txt:Text is widely available for docs; presentations/spreadsheets are less reliable,
            # but this gracefully degrades.
            code, out, err = _run_command([
                str(exe), "--headless", "--convert-to", "txt:Text", "--outdir", tmp, str(path)
            ], timeout=max(DEFAULT_SUBPROCESS_TIMEOUT, 60))
            if code == 0:
                candidates = sorted(Path(tmp).glob("*.txt"))
                for c in candidates:
                    text, extractor, ws = _read_text_file(c)
                    if text.strip():
                        return text, "libreoffice", {}, warnings + ws
            warnings.append(f"libreoffice failed: {(err or out).strip()[:300]}")
    return "", "legacy-office", {}, warnings + ["no legacy office extractor available"]


def _extract_image_metadata(path: Path) -> tuple[str, str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    metadata: dict[str, Any] = {}
    try:
        from PIL import Image  # type: ignore
        with Image.open(path) as img:
            metadata.update({
                "width": img.width,
                "height": img.height,
                "mode": img.mode,
                "format": img.format,
                "frame_count": getattr(img, "n_frames", None),
            })
            try:
                metadata["info_keys"] = sorted(list(img.info.keys()))[:30]
            except Exception:
                pass
    except Exception as e:
        warnings.append(f"PIL metadata unavailable/failed: {e}")
    if path.suffix.lower() in {".tif", ".tiff"}:
        try:
            import tifffile  # type: ignore
            with tifffile.TiffFile(str(path)) as tif:
                metadata.update({
                    "tiff_pages": len(tif.pages),
                    "ome_metadata_present": bool(getattr(tif, "ome_metadata", None)),
                    "series_count": len(tif.series),
                })
                if tif.series:
                    metadata["series_shapes"] = [list(s.shape) for s in tif.series[:5]]
        except Exception as e:
            warnings.append(f"tifffile metadata unavailable/failed: {e}")
    return "", "image-metadata", metadata, warnings


def _extract_pdf_tables_hint(text: str) -> dict[str, Any]:
    # Lightweight table hinting: does not replace Camelot/Tabula, but tells the RAG layer
    # that the document likely contains tabular content worth opening manually.
    lines = [line for line in text.splitlines() if line.strip()]
    tabular_lines = sum(
        1 for line in lines if ("\t" in line or len(re.findall(r"\s{2,}", line)) >= 3)
    )
    return {
        "tabular_line_count_guess": tabular_lines,
        "probably_contains_tables": tabular_lines >= 5,
    }



def _safe_relative_to(path: Path, base_folder: Path | None) -> str:
    if base_folder is None:
        return path.name
    try:
        return str(path.resolve().relative_to(base_folder.resolve()))
    except ValueError:
        return path.name


def _extract_file(path: Path, base_folder: Path | None = None) -> ExtractionResult:
    rel = _safe_relative_to(path, base_folder)
    ext = path.suffix.lower()
    stat = _safe_stat(path)
    mime_type, _ = mimetypes.guess_type(str(path))
    res = ExtractionResult(
        path=rel,
        name=path.name,
        extension=ext,
        document_kind=DOCUMENT_KIND_BY_EXTENSION.get(ext, _classify_asset_type(ext, rel, path.name)),
        mime_type=mime_type,
        size_bytes=stat["size_bytes"],
        modified_at=stat["modified_at"],
        sha256=_sha256_file(path),
    )

    if stat["size_bytes"] > DEFAULT_MAX_FILE_BYTES:
        res.status = "skipped"
        res.warnings.append(f"file exceeds DIGITAL_TWIN_MAX_FILE_BYTES={DEFAULT_MAX_FILE_BYTES}")
        res.extractor = "size-guard"
        return res

    if ext not in SCANNABLE_EXTENSIONS:
        res.status = "skipped"
        res.warnings.append("unsupported extension")
        return res

    if ext in IMAGE_EXTENSIONS:
        text, extractor, metadata, warnings = _extract_image_metadata(path)
        res.status = "metadata_only"
        res.extractor = extractor
        res.metadata.update(metadata)
        res.warnings.extend(warnings)
        return res

    if ext in VIDEO_EXTENSIONS or ext in {".parquet", ".feather", ".h5", ".hdf5", ".rds", ".rdata"}:
        res.status = "metadata_only"
        res.extractor = "metadata-only"
        res.metadata["reason"] = "binary data/video; text extraction intentionally skipped"
        return res

    try:
        text = ""
        extractor = "none"
        metadata: dict[str, Any] = {}
        warnings: list[str] = []

        if ext == ".pdf":
            text, extractor, metadata, warnings = _extract_pdf_text(path)
        elif ext in {".docx", ".dotx"}:
            text, extractor, metadata, warnings = _extract_docx_text(path)
        elif ext in {".odt", ".odp", ".ods"}:
            text, extractor, metadata, warnings = _extract_odt_text(path)
        elif ext == ".pptx":
            text, extractor, metadata, warnings = _extract_pptx_text(path)
        elif ext in {".xlsx"}:
            text, extractor, metadata, warnings = _extract_xlsx_text(path)
        elif ext in {".csv", ".tsv"}:
            text, extractor, metadata, warnings = _extract_csv_text(path, "\t" if ext == ".tsv" else None)
        elif ext == ".ipynb":
            text, extractor, metadata, warnings = _extract_ipynb_text(path)
        elif ext in {".json", ".jsonl"}:
            text, extractor, metadata, warnings = _extract_json_text(path)
        elif ext in {".html", ".htm", ".xml", ".svg", ".drawio"}:
            text, extractor, metadata, warnings = _extract_xmlish_text(path)
        elif ext == ".rtf":
            text, extractor, metadata, warnings = _extract_rtf_text(path)
        elif ext in {".doc", ".ppt", ".xls"}:
            text, extractor, metadata, warnings = _extract_legacy_office_via_cli(path)
        else:
            text, extractor, warnings = _read_text_file(path)
            metadata = {}

        text = _normalize_extracted_text(text)
        text, truncated = _limit_text(text)
        if truncated:
            warnings.append(f"text truncated to {DEFAULT_MAX_TEXT_CHARS_PER_FILE} characters")
        res.extractor = extractor
        res.text = text
        res.metadata.update(metadata)
        if ext == ".pdf" and text:
            res.metadata.update(_extract_pdf_tables_hint(text))
        res.warnings.extend(warnings)
        res.char_count = len(text)
        res.word_count = _count_words(text)
        res.excerpt = _clean_for_excerpt(text)
        res.title = _guess_title_from_text_or_name(text, path.name)
        if ext == ".pdf" and not text.strip():
            res.status = "metadata_only"
            res.metadata["needs_ocr"] = True
            res.metadata["ocr_reason"] = "scanned_or_image_pdf"
            res.warnings.append("PDF has no text layer — eligible for OCR queue")
        elif ext in IMAGE_EXTENSIONS:
            res.status = "metadata_only"
            res.metadata["needs_ocr"] = True
            res.metadata["ocr_reason"] = "image_asset"
        else:
            res.status = "extracted" if text.strip() else "empty"
        res.chunks = _chunk_text(text, rel) if text.strip() else []
        return res
    except Exception as e:  # pragma: no cover
        LOGGER.warning("Extraction failed for %s: %s", path, e)
        res.status = "failed"
        res.errors.append(str(e))
        return res


def _guess_title_from_text_or_name(text: str, name: str) -> str:
    if text:
        for line in text.splitlines()[:40]:
            candidate = _clean(line)
            if 8 <= len(candidate) <= 160 and not _is_junk_title_line(candidate):
                return candidate
    return _humanize_filename(name)


def _classify_file(rel_path: str) -> str:  # type: ignore[override]
    p = rel_path.lower()
    name = Path(rel_path).name.lower()
    ext = Path(rel_path).suffix.lower()
    if "readme" in name or name in {"overview.md", "summary.md"}:
        return "overview"
    if any(k in p for k in ("abstract", "essb", "aacr", "eacr", "manuscript", "poster", "publication", "paper")):
        return "dissemination"
    if any(k in p for k in ("meeting", "logbook", "progress", "update", "minutes")):
        return "activity"
    if any(k in p for k in ("method", "protocol", "experiment", "wet-lab", "wet_lab", "sop")):
        return "protocol"
    if ext in {".py", ".r", ".sh", ".sql", ".ipynb"} or "pipeline" in p or "script" in p:
        return "pipeline"
    if any(k in p for k in ("figure", "figures", "plot", "image", "supplementary")):
        return "figure"
    if ext in DATA_EXTENSIONS:
        return "data"
    if ext in PRESENTATION_EXTENSIONS:
        return "presentation"
    if ext in DOCUMENT_EXTENSIONS:
        return "document"
    return "document"


def _scan_folder(folder: Path) -> list[dict[str, Any]]:  # type: ignore[override]
    files = []
    if not folder or not folder.exists():
        return files
    for path in folder.rglob("*"):
        if not path.is_file():
            continue
        if any(part in SKIP_PARTS for part in path.parts):
            continue
        ext = path.suffix.lower()
        if ext not in SCANNABLE_EXTENSIONS:
            continue
        rel = str(path.relative_to(folder))
        stat = _safe_stat(path)
        files.append({
            "path": rel,
            "name": path.name,
            "extension": ext,
            "size_bytes": stat["size_bytes"],
            "modified_at": stat["modified_at"],
            "category": _classify_file(rel),
            "document_kind": DOCUMENT_KIND_BY_EXTENSION.get(ext, _classify_asset_type(ext, rel, path.name)),
            "folder": str(path.parent.relative_to(folder)) if path.parent != folder else ".",
        })
    return sorted(files, key=lambda f: f["path"])


def _scan_all_assets(folder: Path) -> list[dict[str, Any]]:  # type: ignore[override]
    assets: list[dict[str, Any]] = []
    if not folder or not folder.exists():
        return assets
    for path in folder.rglob("*"):
        if not path.is_file():
            continue
        if any(part in SKIP_PARTS for part in path.parts):
            continue
        ext = path.suffix.lower()
        if ext not in SCANNABLE_EXTENSIONS:
            continue
        rel = str(path.relative_to(folder))
        stat = _safe_stat(path)
        name = path.name
        asset_type = _classify_asset_type(ext, rel, name)
        lifecycle = _classify_lifecycle_section(rel)
        priority = _figure_priority(rel, name) if asset_type in ("figure", "image") else 0
        previewable = ext in IMAGE_EXTENSIONS or ext in {".pdf", ".pptx", ".docx", ".xlsx", ".html", ".svg"}
        assets.append({
            "path": rel,
            "name": name,
            "extension": ext,
            "size_bytes": stat["size_bytes"],
            "modified_at": stat["modified_at"],
            "asset_type": asset_type,
            "document_kind": DOCUMENT_KIND_BY_EXTENSION.get(ext, asset_type),
            "lifecycle_section": lifecycle,
            "category": _classify_file(rel),
            "folder": str(path.parent.relative_to(folder)) if path.parent != folder else ".",
            "priority": priority,
            "previewable": previewable,
            "excerpt": "",
        })
    return sorted(assets, key=lambda a: a["path"])


def _extract_project_documents(folder: Path, assets: list[dict[str, Any]]) -> tuple[list[ExtractionResult], list[dict[str, Any]], dict[str, Any]]:
    records: list[ExtractionResult] = []
    chunks: list[dict[str, Any]] = []
    status_counts: dict[str, int] = {}
    extractor_counts: dict[str, int] = {}
    extension_counts: dict[str, int] = {}
    errors: list[dict[str, str]] = []

    for asset in assets:
        ext = asset.get("extension", "")
        extension_counts[ext] = extension_counts.get(ext, 0) + 1
        path = folder / asset["path"]
        if ext not in SCANNABLE_EXTENSIONS:
            continue
        result = _extract_file(path, folder)
        records.append(result)
        status_counts[result.status] = status_counts.get(result.status, 0) + 1
        extractor_counts[result.extractor] = extractor_counts.get(result.extractor, 0) + 1
        if result.errors:
            errors.append({"path": result.path, "error": "; ".join(result.errors)[:500]})
        if result.chunks:
            chunks.extend(result.chunks)
        # Copy useful extraction summary into asset for UI/library.
        asset["excerpt"] = result.excerpt
        asset["title"] = result.title
        asset["extraction_status"] = result.status
        asset["extractor"] = result.extractor
        asset["word_count"] = result.word_count
        asset["char_count"] = result.char_count
        asset["sha256"] = result.sha256
        asset["metadata"] = result.metadata

    summary = {
        "total_scannable_assets": len(assets),
        "extracted_records": len(records),
        "chunk_count": len(chunks),
        "status_counts": status_counts,
        "extractor_counts": extractor_counts,
        "extension_counts": extension_counts,
        "errors": errors[:50],
        "limits": {
            "max_text_chars_per_file": DEFAULT_MAX_TEXT_CHARS_PER_FILE,
            "chunk_chars": DEFAULT_CHUNK_CHARS,
            "chunk_overlap": DEFAULT_CHUNK_OVERLAP,
            "max_pdf_pages": DEFAULT_MAX_PDF_PAGES,
            "max_pptx_slides": DEFAULT_MAX_PPTX_SLIDES,
            "max_xlsx_sheets": DEFAULT_MAX_XLSX_SHEETS,
            "max_xlsx_rows": DEFAULT_MAX_XLSX_ROWS,
            "max_file_bytes": DEFAULT_MAX_FILE_BYTES,
        },
        "optional_extractors_detected": {
            "pypdf": _python_module_exists("pypdf"),
            "PyPDF2": _python_module_exists("PyPDF2"),
            "python_docx": _python_module_exists("docx"),
            "python_pptx": _python_module_exists("pptx"),
            "openpyxl": _python_module_exists("openpyxl"),
            "beautifulsoup4": _python_module_exists("bs4"),
            "Pillow": _python_module_exists("PIL"),
            "tifffile": _python_module_exists("tifffile"),
            "pdftotext_cli": _command_exists("pdftotext"),
            "libreoffice_cli": _command_exists("libreoffice") or _command_exists("soffice"),
            "textutil_cli": _command_exists("textutil"),
        },
    }
    return records, chunks, summary


def _python_module_exists(module_name: str) -> bool:
    try:
        __import__(module_name)
        return True
    except Exception:
        return False


def _build_content_library(assets: list[dict[str, Any]]) -> dict[str, Any]:  # type: ignore[override]
    section_defs = {
        r[0]: {
            "id": r[0], "label": r[1], "figures": [], "documents": [], "presentations": [],
            "data_files": [], "text_files": [], "videos": [], "code_files": [], "counts": {}
        }
        for r in LIFECYCLE_RULES
    }
    figures_pool: list[dict] = []
    totals: dict[str, int] = {}

    for a in assets:
        sec_id = a.get("lifecycle_section", "root")
        sec = section_defs.get(sec_id, section_defs["root"])
        atype = a.get("asset_type", "other")
        ext = a.get("extension", "")
        if ext in {".py", ".r", ".sh", ".sql", ".ipynb", ".js", ".ts", ".tsx", ".jsx", ".go"}:
            atype_bucket = "code"
        else:
            atype_bucket = atype
        totals[atype_bucket] = totals.get(atype_bucket, 0) + 1
        totals["all"] = totals.get("all", 0) + 1
        item = {
            "path": a["path"],
            "name": a["name"],
            "extension": ext,
            "size_bytes": a["size_bytes"],
            "modified_at": a.get("modified_at"),
            "asset_type": atype,
            "document_kind": a.get("document_kind"),
            "previewable": a.get("previewable", False),
            "excerpt": a.get("excerpt", ""),
            "title": a.get("title", ""),
            "extraction_status": a.get("extraction_status"),
            "extractor": a.get("extractor"),
            "word_count": a.get("word_count", 0),
        }
        bucket_key = {
            "figure": "figures", "image": "figures", "document": "documents",
            "presentation": "presentations", "data": "data_files", "text": "text_files",
            "video": "videos", "code": "code_files",
        }.get(atype_bucket, "documents")
        sec["counts"][atype_bucket] = sec["counts"].get(atype_bucket, 0) + 1
        if atype in ("figure", "image"):
            figures_pool.append({**item, "section": sec_id, "section_label": sec["label"], "priority": a.get("priority", 0)})
        lst = sec[bucket_key]
        if len(lst) < 80:
            lst.append(item)

    figures_gallery = sorted(figures_pool, key=lambda x: (-x.get("priority", 0), x["path"]))[:96]
    sections = []
    for sid in [r[0] for r in LIFECYCLE_RULES]:
        sec = section_defs[sid]
        if not sec["counts"]:
            continue
        sections.append({
            "id": sec["id"],
            "label": sec["label"],
            "figures": sec["figures"],
            "documents": sec["documents"],
            "presentations": sec["presentations"],
            "data_files": sec["data_files"],
            "text_files": sec["text_files"],
            "videos": sec["videos"],
            "code_files": sec["code_files"],
            "counts": sec["counts"],
            "total_files": sum(sec["counts"].values()),
        })

    return {
        "sections": sections,
        "figures_gallery": figures_gallery,
        "totals": totals,
        "figure_count": totals.get("figure", 0) + totals.get("image", 0),
    }


def _combine_text_records(records: list[ExtractionResult], *, max_total_chars: int = 2_000_000) -> str:
    pieces: list[str] = []
    total = 0
    for r in records:
        if not r.text:
            continue
        header = f"\n\n# Source: {r.path}\n"
        part = header + r.text
        if total + len(part) > max_total_chars:
            pieces.append(part[: max(0, max_total_chars - total)])
            break
        pieces.append(part)
        total += len(part)
    return "\n".join(pieces)


def _parse_dissemination_from_records(records: list[ExtractionResult]) -> list[dict[str, Any]]:
    items = []
    for r in records:
        path_lower = r.path.lower()
        if not any(k in path_lower for k in ("abstract", "essb", "aacr", "eacr", "manuscript", "poster", "paper", "publication")):
            continue
        title = r.title or r.name.replace("_", " ").rsplit(".", 1)[0]
        conference = ""
        for conf in ("AACR", "EACR", "EMBO", "ESSB", "SCSO", "ESMO", "ASCO"):
            if conf.lower() in path_lower:
                conference = conf
        year_m = re.search(r"20\d{2}", r.path)
        dois = DOI_PATTERN.findall(r.text or "")
        items.append({
            "title": _clean(title),
            "author": "",
            "conference": conference,
            "year": int(year_m.group()) if year_m else None,
            "doi": dois[0] if dois else None,
            "type": "abstract" if "abstract" in path_lower else ("poster" if "poster" in path_lower else "manuscript"),
            "source_file": r.path,
        })
    return items[:80]


def _parse_protocols_from_records(records: list[ExtractionResult]) -> list[dict[str, Any]]:
    protocols = []
    for r in records:
        path_lower = r.path.lower()
        if not any(k in path_lower for k in ("protocol", "method", "sop", "experiment", "wet-lab", "wet_lab")):
            continue
        content = r.text or ""
        steps = []
        for line in content.splitlines():
            line = line.strip()
            if re.match(r"^\d+[\.)]\s+", line) or line.startswith("- ") or line.startswith("* "):
                step = _clean(re.sub(r"^\d+[\.)]\s+|^[-*]\s+", "", line))
                if 10 < len(step) < 350:
                    steps.append(step)
        if not steps:
            # Fallback: first good sentences.
            sentences = re.split(r"(?<=[.!?])\s+", _clean(content))
            steps = [s for s in sentences if 20 < len(s) < 350][:12]
        protocols.append({
            "title": r.title or r.name.rsplit(".", 1)[0].replace("_", " "),
            "category": str(Path(r.path).parent),
            "steps": steps[:20],
            "source_file": r.path,
        })
    return protocols[:60]


def _extract_action_items(text: str, source_file: str = "combined") -> list[dict[str, Any]]:
    items = []
    for m in TODO_ITEM.finditer(text or ""):
        body = _clean(m.group(1))
        if body:
            items.append({"title": body[:160], "summary": body, "source_file": source_file, "category": "action"})
    # Also catch checkbox tasks.
    for m in re.finditer(r"(?:^|\n)\s*[-*]\s*\[[ xX]?\]\s+(.+)", text or ""):
        body = _clean(m.group(1))
        if body:
            items.append({"title": body[:160], "summary": body, "source_file": source_file, "category": "action"})
    return items[:150]


def _parse_methods_and_findings(text: str) -> dict[str, list[str]]:
    methods: list[str] = []
    findings: list[str] = []
    risks: list[str] = []
    for line in text.splitlines():
        clean = _clean(line)
        low = clean.lower()
        if 20 <= len(clean) <= 300:
            if any(k in low for k in ("method", "workflow", "pipeline", "analysis", "segmentation", "quantification", "deconvolution", "stitching")):
                methods.append(clean)
            if any(k in low for k in ("finding", "identified", "showed", "correlat", "significant", "suggest", "revealed")):
                findings.append(clean)
            if any(k in low for k in ("todo", "risk", "missing", "failed", "warning", "problem", "issue", "limitation")):
                risks.append(clean)
    return {
        "methods": _dedupe_keep_order(methods)[:80],
        "findings": _dedupe_keep_order(findings)[:80],
        "risks_or_open_items": _dedupe_keep_order(risks)[:80],
    }


def _dedupe_keep_order(items: list[str]) -> list[str]:
    seen = set()
    out = []
    for item in items:
        key = re.sub(r"\s+", " ", item.lower()).strip()
        if key and key not in seen:
            seen.add(key)
            out.append(item)
    return out


def _guess_project_summary(text: str) -> str:
    candidates = []
    for pattern in (
        r"(?:Core Objective|Objective|Aim|Description)\s*[:\n]+(.{80,900}?)(?:\n\n|\n#|$)",
        r"(?:This project|The project|We aim|Aim is)\s+(.{80,700}?)(?:\n\n|\n#|$)",
    ):
        m = re.search(pattern, text or "", re.I | re.S)
        if m:
            candidates.append(_clean(m.group(1)))
    if candidates:
        return candidates[0][:900]
    return _clean(text or "")[:900]


# --- Raw Knowledge Vault ingestion policies (metadata-first, no forced categorization) ---

VAULT_LARGE_BINARY_SUFFIXES = (
    ".ome.tiff", ".ome.tif", ".tif", ".tiff", ".h5", ".hdf5", ".rds", ".rdata", ".ome.xml",
)
VAULT_SCRIPT_EXTENSIONS = {".py", ".r", ".rmd", ".ipynb", ".sh", ".jl", ".m"}
VAULT_LOG_EXTENSIONS = {".log", ".out", ".err"}
VAULT_LOG_MAX_LINES = int(os.environ.get("VAULT_LOG_MAX_LINES", "5000"))
VAULT_LOG_ERROR_PATTERNS = re.compile(
    r"(?i)(error|exception|traceback|fatal|failed|errno|warning:)",
)
VAULT_EXCEL_SAMPLE_ROWS = int(os.environ.get("VAULT_EXCEL_SAMPLE_ROWS", "5"))


def _path_suffix_lower(path: Path) -> str:
    name = path.name.lower()
    for suffix in sorted(VAULT_LARGE_BINARY_SUFFIXES, key=len, reverse=True):
        if name.endswith(suffix):
            return suffix
    return path.suffix.lower()


def is_vault_large_binary(path: Path) -> bool:
    return _path_suffix_lower(path) in VAULT_LARGE_BINARY_SUFFIXES or path.suffix.lower() in {
        ".h5", ".hdf5", ".rds", ".rdata",
    }


def _extract_excel_vault_metadata(path: Path) -> tuple[dict[str, Any], list[str]]:
    """Sheet names, column headers, and limited sample rows for vault metadata_json."""
    warnings: list[str] = []
    metadata: dict[str, Any] = {"vault_policy": "excel_schema_preview"}
    try:
        from openpyxl import load_workbook  # type: ignore
        wb = load_workbook(filename=str(path), read_only=True, data_only=True)
        sheets_out: list[dict[str, Any]] = []
        for sheet_name in wb.sheetnames[:DEFAULT_MAX_XLSX_SHEETS]:
            ws = wb[sheet_name]
            headers: list[str] = []
            samples: list[list[str]] = []
            for i, row in enumerate(ws.iter_rows(max_row=DEFAULT_MAX_XLSX_ROWS + 1, values_only=True)):
                vals = [str(v).strip() if v is not None else "" for v in row]
                if i == 0:
                    headers = [v for v in vals if v]
                elif len(samples) < VAULT_EXCEL_SAMPLE_ROWS:
                    if any(vals):
                        samples.append(vals[: min(32, len(vals))])
            sheets_out.append({
                "name": sheet_name,
                "column_headers": headers[:64],
                "sample_rows": samples,
                "column_count": len(headers) or max((len(r) for r in samples), default=0),
            })
        metadata["sheet_count"] = len(wb.sheetnames)
        metadata["sheets"] = sheets_out
        try:
            wb.close()
        except Exception:
            pass
        return metadata, warnings
    except Exception as e:
        warnings.append(f"openpyxl metadata failed: {e}")
        if path.suffix.lower() in {".csv", ".tsv"}:
            text, _, md, ws = _extract_csv_text(path)
            metadata.update(md)
            metadata["preview_text_chars"] = len(text)
            return metadata, warnings + ws
        return metadata, warnings


def _extract_script_vault_summary(path: Path) -> tuple[dict[str, Any], list[str]]:
    warnings: list[str] = []
    ext = path.suffix.lower()
    metadata: dict[str, Any] = {"vault_policy": "script_summary"}
    if ext == ".ipynb":
        try:
            nb = json.loads(path.read_text(encoding="utf-8", errors="replace"))
            cells = nb.get("cells", [])
            code_lines = 0
            imports: list[str] = []
            funcs: list[str] = []
            for cell in cells:
                if cell.get("cell_type") != "code":
                    continue
                src = cell.get("source", "")
                if isinstance(src, list):
                    src = "".join(src)
                code_lines += src.count("\n") + (1 if src.strip() else 0)
                for m in re.finditer(r"^\s*(?:import|from)\s+([\w.]+)", src, re.M):
                    imports.append(m.group(1))
                for m in re.finditer(r"^\s*def\s+(\w+)", src, re.M):
                    funcs.append(m.group(1))
            metadata.update({
                "cell_count": len(cells),
                "code_line_count": code_lines,
                "imports": _dedupe_keep_order(imports)[:40],
                "functions": _dedupe_keep_order(funcs)[:40],
                "kernel": (nb.get("metadata", {}).get("kernelspec") or {}).get("display_name"),
            })
            return metadata, warnings
        except Exception as e:
            return metadata, [str(e)]

    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except OSError as e:
        return metadata, [str(e)]
    lines = text.splitlines()
    metadata["line_count"] = len(lines)
    imports: list[str] = []
    funcs: list[str] = []
    if ext == ".py":
        for m in re.finditer(r"^\s*(?:import|from)\s+([\w.]+)", text, re.M):
            imports.append(m.group(1))
        for m in re.finditer(r"^\s*def\s+(\w+)", text, re.M):
            funcs.append(m.group(1))
        for m in re.finditer(r"^\s*class\s+(\w+)", text, re.M):
            funcs.append(f"class:{m.group(1)}")
    elif ext == ".r":
        for m in re.finditer(r"^\s*(?:library|require)\s*\(\s*['\"]?([\w.]+)", text, re.M):
            imports.append(m.group(1))
        for m in re.finditer(r"^\s*(\w+)\s*<-\s*function", text, re.M):
            funcs.append(m.group(1))
    metadata["imports"] = _dedupe_keep_order(imports)[:40]
    metadata["functions"] = _dedupe_keep_order(funcs)[:40]
    if len(text) > 8000:
        metadata["text_preview"] = text[:8000]
        warnings.append("script text truncated in metadata preview")
    else:
        metadata["text_preview"] = text
    return metadata, warnings


def _extract_log_vault_summary(path: Path) -> tuple[dict[str, Any], list[str]]:
    warnings: list[str] = []
    metadata: dict[str, Any] = {"vault_policy": "log_summary"}
    try:
        with path.open("r", encoding="utf-8", errors="replace") as fh:
            lines: list[str] = []
            error_lines: list[str] = []
            for i, line in enumerate(fh):
                if i >= VAULT_LOG_MAX_LINES:
                    warnings.append(f"log read capped at {VAULT_LOG_MAX_LINES} lines")
                    break
                lines.append(line.rstrip("\n"))
                if VAULT_LOG_ERROR_PATTERNS.search(line):
                    error_lines.append(line.strip()[:500])
        metadata["line_count"] = len(lines)
        metadata["error_line_count"] = len(error_lines)
        metadata["error_lines_sample"] = error_lines[:50]
        metadata["tail_sample"] = lines[-20:] if lines else []
        return metadata, warnings
    except OSError as e:
        return metadata, [str(e)]


def extract_for_vault(path: Path, base_folder: Path | None = None) -> ExtractionResult:
    """Vault-oriented extraction: always returns a result; large binaries are metadata-only."""
    rel = _safe_relative_to(path, base_folder)
    ext = path.suffix.lower()
    suffix = _path_suffix_lower(path)
    stat = _safe_stat(path)
    mime_type, _ = mimetypes.guess_type(str(path))
    res = ExtractionResult(
        path=rel,
        name=path.name,
        extension=ext or suffix,
        document_kind=DOCUMENT_KIND_BY_EXTENSION.get(ext, _classify_asset_type(ext, rel, path.name)),
        mime_type=mime_type,
        size_bytes=stat["size_bytes"],
        modified_at=stat["modified_at"],
        sha256=_sha256_file(path),
    )

    if stat["size_bytes"] > DEFAULT_MAX_FILE_BYTES:
        res.status = "skipped"
        res.warnings.append(f"file exceeds DIGITAL_TWIN_MAX_FILE_BYTES={DEFAULT_MAX_FILE_BYTES}")
        res.extractor = "size-guard"
        res.metadata["vault_policy"] = "skipped_oversize"
        return res

    if is_vault_large_binary(path):
        res.status = "metadata_only"
        res.extractor = "vault-large-binary"
        res.metadata.update({
            "vault_policy": "large_binary_metadata_only",
            "logical_path": rel,
            "size_bytes": stat["size_bytes"],
            "mime_type": mime_type,
            "checksum_sha256": res.sha256,
            "suffix": suffix,
        })
        if suffix in {".tif", ".tiff"} or "ome" in suffix:
            text, extractor, md, warnings = _extract_image_metadata(path)
            res.metadata.update(md)
            res.warnings.extend(warnings)
            res.extractor = extractor or res.extractor
        return res

    if ext in {".xlsx", ".xls"} or path.suffix.lower() in {".csv", ".tsv"}:
        md, warnings = _extract_excel_vault_metadata(path)
        res.metadata.update(md)
        res.warnings.extend(warnings)
        if ext == ".xlsx":
            text, extractor, xmd, xw = _extract_xlsx_text(path)
            res.text = _normalize_extracted_text(text)
            res.extractor = extractor
            res.metadata.update(xmd)
            res.warnings.extend(xw)
        elif path.suffix.lower() in {".csv", ".tsv"}:
            text, extractor, cmd, cw = _extract_csv_text(path)
            res.text = _normalize_extracted_text(text)
            res.extractor = extractor
            res.metadata.update(cmd)
            res.warnings.extend(cw)
        res.status = "extracted" if res.text.strip() else "metadata_only"
        res.char_count = len(res.text)
        res.excerpt = _clean_for_excerpt(res.text)
        return res

    if ext in VAULT_SCRIPT_EXTENSIONS:
        md, warnings = _extract_script_vault_summary(path)
        res.metadata.update(md)
        res.warnings.extend(warnings)
        res.status = "metadata_only"
        res.extractor = "vault-script-summary"
        preview = (md.get("text_preview") or "")[:4000]
        if preview:
            res.text = preview
            res.excerpt = _clean_for_excerpt(preview)
            res.char_count = len(preview)
        return res

    if ext in VAULT_LOG_EXTENSIONS:
        md, warnings = _extract_log_vault_summary(path)
        res.metadata.update(md)
        res.warnings.extend(warnings)
        res.status = "metadata_only"
        res.extractor = "vault-log-summary"
        return res

    try:
        inner = _extract_file(path, base_folder)
        res.status = inner.status
        res.extractor = inner.extractor
        res.text = inner.text
        res.excerpt = inner.excerpt
        res.title = inner.title
        res.metadata.update(inner.metadata)
        res.metadata["vault_policy"] = res.metadata.get("vault_policy") or "standard_extract"
        res.warnings.extend(inner.warnings)
        res.errors.extend(inner.errors)
        res.char_count = inner.char_count
        res.word_count = inner.word_count
        res.chunks = inner.chunks
        if inner.status == "failed":
            res.metadata["error"] = "; ".join(inner.errors)[:2000]
        return res
    except Exception as e:
        res.status = "failed"
        res.errors.append(str(e))
        res.metadata["error"] = str(e)
        res.metadata["vault_policy"] = "extract_exception"
        return res


def vault_extraction_status(result: ExtractionResult) -> str:
    """Map ExtractionResult.status to platform.raw_asset_vault.extraction_status."""
    mapping = {
        "extracted": "extracted",
        "metadata_only": "metadata_only",
        "empty": "empty",
        "skipped": "skipped",
        "failed": "failed",
        "not_attempted": "not_started",
    }
    return mapping.get(result.status, result.status)
