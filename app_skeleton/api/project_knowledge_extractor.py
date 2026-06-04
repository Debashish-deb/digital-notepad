import os
import json
import uuid
import hashlib
import psycopg
from pathlib import Path
from app_skeleton.api.supabase_config import postgres_conn
from app_skeleton.api.paths import PROJECTS_ROOT

def generate_document_code(project_code: str, relative_path: str) -> str:
    hash_object = hashlib.md5(f"{project_code}:{relative_path}".encode())
    return f"{project_code.upper()}-DOC-{hash_object.hexdigest()[:8].upper()}"

def generate_chunk_uid(doc_code: str, chunk_index: int) -> str:
    return f"{doc_code}-CHK-{chunk_index}"

import re
def calculate_confidence(text: str) -> float:
    if not text:
        return 0.0
    text = text.strip()
    if not text:
        return 0.0
    tokens = text.split()
    if not tokens:
        return 0.0
    valid_tokens = [t for t in tokens if re.match(r'^[a-zA-Z0-9\-\.,\?!\(\)\[\]{};:\'"/\\+=$%&@<>_]+$', t)]
    return (len(valid_tokens) / len(tokens)) * 100.0


def chunk_text(text: str, max_tokens: int = 1000) -> list[str]:
    max_chars = max_tokens * 4
    paragraphs = text.split('\n\n')
    chunks = []
    current_chunk = ""
    for para in paragraphs:
        # Force split massive paragraphs without double newlines
        while len(para) > max_chars:
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
                current_chunk = ""
            chunks.append(para[:max_chars].strip())
            para = para[max_chars:]
            
        if len(current_chunk) + len(para) > max_chars:
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            current_chunk = para + "\n\n"
        else:
            current_chunk += para + "\n\n"
    if current_chunk.strip():
        chunks.append(current_chunk.strip())
    if not chunks and text.strip():
        chunks.append(text.strip()[:max_chars])
    return chunks

def extract_and_ingest_project(project_code: str) -> dict:
    project_dir = PROJECTS_ROOT / project_code
    if not project_dir.is_dir():
        # Fallback 1: check if it's in the repo root directly if PROJECTS_ROOT failed
        fallback_dir = project_dir.parent.parent / "projects" / project_code
        if fallback_dir.is_dir():
            project_dir = fallback_dir
        else:
            # Fallback 2: fuzzy match in PROJECTS_ROOT (e.g. 4_CellCycle for CellCycle)
            found = False
            for child in PROJECTS_ROOT.iterdir():
                if child.is_dir() and project_code.lower() in child.name.lower():
                    project_dir = child
                    found = True
                    break
            
            if not found:
                raise FileNotFoundError(f"Project directory not found for {project_code}")
    
    valid_extensions = ['.md', '.txt', '.json', '.csv', '.tsv', '.doc', '.docx', '.pdf', '.pptx']
    
    document_sources = []
    document_chunks = []
    
    # Try importing docx, but fail gracefully if not installed
    try:
        import docx
    except ImportError:
        docx = None
        
    try:
        import fitz
    except ImportError:
        fitz = None
        
    try:
        from pptx import Presentation
    except ImportError:
        Presentation = None
        
    try:
        import pypdf
    except ImportError:
        pypdf = None
    
    for dirpath, _, filenames in os.walk(project_dir):
        for filename in filenames:
            ext = os.path.splitext(filename)[1].lower()
            if any(filename.endswith(valid_ext) for valid_ext in valid_extensions):
                if filename in ['extract_script.py', 'cellcycle_extracted_knowledge.json', 'push_to_supabase.py']:
                    continue
                    
                file_path = os.path.join(dirpath, filename)
                relative_path = os.path.relpath(file_path, project_dir)
                
                try:
                    if os.path.getsize(file_path) > 5 * 1024 * 1024:
                        print(f"Skipping {file_path} (exceeds 5MB limit)")
                        continue
                        
                    if ext in ['.doc', '.docx']:
                        if docx:
                            doc_obj = docx.Document(file_path)
                            content = "\n\n".join([p.text for p in doc_obj.paragraphs if p.text.strip()])
                        else:
                            content = ""
                    elif ext == '.pdf':
                        content = ""
                        if fitz:
                            try:
                                with fitz.open(file_path) as doc_obj:
                                    content = "\n".join(page.get_text() for page in doc_obj)
                            except Exception as ex:
                                print(f"PyMuPDF failed for {file_path}: {ex}")
                        
                        if calculate_confidence(content) < 98.0 and pypdf:
                            try:
                                reader = pypdf.PdfReader(file_path)
                                content = "\n".join(page.extract_text() or "" for page in reader.pages)
                            except Exception as ex:
                                print(f"pypdf failed for {file_path}: {ex}")
                                
                        if calculate_confidence(content) < 98.0:
                            print(f"Skipping semantic chunking for {file_path} (confidence < 98%). Treating as visual asset.")
                            content = ""
                    elif ext == '.pptx':
                        content = ""
                        if Presentation:
                            try:
                                prs = Presentation(file_path)
                                text_runs = []
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if hasattr(shape, "text"):
                                            text_runs.append(shape.text)
                                content = "\n\n".join(text_runs)
                            except Exception as ex:
                                print(f"python-pptx failed for {file_path}: {ex}")
                        
                        if calculate_confidence(content) < 98.0:
                            print(f"Skipping semantic chunking for {file_path} (confidence < 98%). Treating as visual asset.")
                            content = ""
                    else:
                        with open(file_path, 'r', encoding='utf-8') as f:
                            content = f.read()
                except Exception as e:
                    print(f"Failed to read {file_path}: {e}")
                    continue
                
                doc_code = generate_document_code(project_code, relative_path)
                source_type = filename.split('.')[-1]
                
                doc_source = {
                    "document_code": doc_code,
                    "title": filename,
                    "source_type": source_type,
                    "metadata": {
                        "project_code": project_code,
                        "relative_path": relative_path,
                        "corpus": "project_workspace",
                        "folder_path": os.path.relpath(dirpath, project_dir)
                    }
                }
                document_sources.append(doc_source)
                
                chunks = chunk_text(content)
                for idx, chunk in enumerate(chunks):
                    doc_chunk = {
                        "document_code": doc_code,
                        "chunk_index": idx,
                        "chunk_uid": generate_chunk_uid(doc_code, idx),
                        "chunk_text": chunk,
                        "token_count": len(chunk) // 4,
                        "metadata": {
                            "char_count": len(chunk)
                        }
                    }
                    document_chunks.append(doc_chunk)

    if not document_sources:
        return {"extracted_docs": 0, "extracted_chunks": 0, "message": "No valid documents found."}

    # Push to database
    conn_uri = postgres_conn()
    if not conn_uri:
        raise ConnectionError("PostgreSQL connection URI not configured.")

    with psycopg.connect(conn_uri) as conn:
        with conn.cursor() as cur:
            # 1. Insert into rag.document_source
            for doc in document_sources:
                cur.execute("""
                    INSERT INTO rag.document_source (document_code, title, source_type, metadata)
                    VALUES (%s, %s, %s, %s)
                    ON CONFLICT (document_code) DO UPDATE
                    SET title = EXCLUDED.title,
                        source_type = EXCLUDED.source_type,
                        metadata = EXCLUDED.metadata
                    RETURNING document_id;
                """, (doc["document_code"], doc["title"], doc["source_type"], json.dumps(doc["metadata"])))
                doc_id = cur.fetchone()[0]
                doc["_db_id"] = doc_id
            
            # Map document_code to generated document_id
            doc_id_map = {doc["document_code"]: doc["_db_id"] for doc in document_sources}
            
            # 2. Insert into rag.document_chunk in batches
            batch_size = 50
            for i, chunk in enumerate(document_chunks):
                doc_id = doc_id_map.get(chunk["document_code"])
                if not doc_id:
                    continue
                cur.execute("""
                    INSERT INTO rag.document_chunk (document_id, chunk_index, chunk_uid, chunk_text, token_count, metadata)
                    VALUES (%s, %s, %s, %s, %s, %s)
                    ON CONFLICT (chunk_uid) DO UPDATE
                    SET chunk_text = EXCLUDED.chunk_text,
                        token_count = EXCLUDED.token_count,
                        metadata = EXCLUDED.metadata;
                """, (doc_id, chunk["chunk_index"], chunk["chunk_uid"], chunk["chunk_text"], chunk["token_count"], json.dumps(chunk["metadata"])))
                
                # Commit periodically to avoid exceeding transaction limits on massive files
                if (i + 1) % batch_size == 0:
                    conn.commit()
            
            # Final commit for remaining chunks
            conn.commit()
    
    return {
        "extracted_docs": len(document_sources),
        "extracted_chunks": len(document_chunks),
        "message": f"Successfully ingested {len(document_sources)} docs."
    }
